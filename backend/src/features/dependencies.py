"""
file_tree_builder.py
=====================

Универсальный построитель дерева из файла любого типа:
.txt / .md, .docx, .pptx, .pdf, .json, .xml / .html

Идея архитектуры (важно для понимания "почему это хорошо"):

1. Единая модель узла (TreeNode) — не зависит от формата файла.
   Все парсеры формата просто ЗАПОЛНЯЮТ одну и ту же структуру.
   Это классический паттерн "Adapter": формат -> единый интерфейс.
   Благодаря этому дерево из docx, pptx и json можно обрабатывать
   ОДНИМ И ТЕМ ЖЕ кодом (печать, экспорт в JSON, поиск, обход и т.д.)

2. Диспетчер по расширению файла (build_tree) выбирает нужный парсер.
   Добавить новый формат = дописать одну функцию parse_XXX и
   зарегистрировать её в словаре PARSERS. Ничего другого менять не надо
   (открытость/закрытость — Open/Closed Principle).

3. docx/pptx/pdf по умолчанию идут через двухшаговый конвейер:
       файл -> markitdown -> нормализованный markdown -> дерево
   Вместо трёх разных парсеров под каждый формат используется ОДНО
   markdown-ядро (_markdown_text_to_tree). markitdown (библиотека
   Microsoft) сам разбирается со спецификой формата — стилями Word,
   слайдами pptx, таблицами — и отдаёт понятный markdown; наш код
   отвечает только за превращение markdown в дерево. Это резко снижает
   число мест, где может быть баг, и работает одинаково хорошо на
   "грязных" файлах, где нет явной XML-разметки структуры.
   json: сама структура dict/list рекурсивно превращается в дерево.
   xml/html: дерево тегов ставится в соответствие дереву узлов 1:к:1.
   Прямые парсеры через python-docx/python-pptx/pdfplumber тоже
   оставлены (parse_*_native) — для случаев, когда нужна 100%-точная
   структура по стилям/таблицам, а не markdown-нормализация.

4. Производительность:
   - Никаких регулярных выражений там, где есть API уровня объектов
     (python-docx/pptx уже парсят XML через lxml — это C-расширение,
     на порядок быстрее ручного разбора строк).
   - Обход файлов делается за один проход O(n), где n — число
     структурных единиц (абзацев, слайдов, строк), а не через
     повторные пересканирования.
   - Для pdf используется постраничная обработка (генератор),
     поэтому большие файлы не парсятся "все и сразу" сверх необходимого.
"""

from __future__ import annotations

from typing import Any

import json
import os
from collections.abc import Callable
from dataclasses import dataclass, field

# ---------------------------------------------------------------------------
# 1. УНИВЕРСАЛЬНАЯ МОДЕЛЬ УЗЛА ДЕРЕВА
# ---------------------------------------------------------------------------


@dataclass
class TreeNode:
    """
    Узел дерева, одинаковый для любого формата файла.

    name      — короткая подпись узла (заголовок, "Slide 1", тег и т.п.)
    node_type — тип узла ('document', 'heading', 'paragraph', 'slide',
                'shape', 'page', 'line', 'json_object', 'json_array',
                'xml_tag', ...) — нужен, чтобы разные обработчики могли
                по-разному реагировать на разные типы узлов.
    level     — глубина вложенности (0 = корень). Дублирует глубину в
                дереве, но хранить явно удобно для быстрой фильтрации
                "показать только заголовки уровня <= 2" без обхода дерева.
    content   — сам текст/значение узла (может быть None у "контейнеров")
    children  — список дочерних TreeNode
    """

    name: str
    node_type: str = "node"
    level: int = 0
    content: str | None = None
    children: list[TreeNode] = field(default_factory=list)

    def add_child(self, child: TreeNode) -> TreeNode:
        self.children.append(child)
        return child

    def to_dict(self) -> dict:
        """Сериализация в обычный dict — пригодится для json.dump,
        для передачи во фронтенд, для сохранения в БД и т.д."""
        return {
            "name": self.name,
            "type": self.node_type,
            "level": self.level,
            "content": self.content,
            "children": [c.to_dict() for c in self.children],
        }

    def print_tree(self, prefix: str = "", is_last: bool = True) -> None:
        """Красивая ASCII-печать дерева (как `tree` в консоли)."""
        connector = "└── " if is_last else "├── "
        label = self.name
        if self.content:
            snippet = self.content.strip().replace("\n", " ")
            if len(snippet) > 60:
                snippet = snippet[:57] + "..."
            if snippet:
                label += f"  «{snippet}»"
        print(prefix + connector + label)
        new_prefix = prefix + ("    " if is_last else "│   ")
        for i, child in enumerate(self.children):
            child.print_tree(new_prefix, i == len(self.children) - 1)

    def count_nodes(self) -> int:
        return 1 + sum(c.count_nodes() for c in self.children)


# ---------------------------------------------------------------------------
# 2. ПАРСЕРЫ ФОРМАТОВ — каждый возвращает готовый TreeNode (корень)
# ---------------------------------------------------------------------------

import pathlib
import re

_SLIDE_MARKER_RE = re.compile(r"^<!--\s*Slide number:\s*(\d+)\s*-->$")
_IMAGE_LINE_RE = re.compile(r"^!\[.*\]\(.*\)$")


def _markdown_text_to_tree(text: str, root_name: str, skip_images: bool = True) -> TreeNode:
    """
    Общее ядро для построения дерева из markdown-текста.
    Используется и для родных .md/.txt файлов, и для markdown,
    который получен из docx/pptx/pdf через markitdown (см. ниже) —
    поэтому вся логика распознавания структуры живёт в ОДНОМ месте,
    а не дублируется в каждом парсере формата.

    Распознаёт:
    - '# ... ######' — markdown-заголовки -> уровень = число '#'
    - '<!-- Slide number: N -->' — маркер слайда, которым markitdown
      размечает границы слайдов в pptx -> отдельный узел-заголовок
    - '- ' / '* ' — пункты списка -> вложены под текущий заголовок,
      дополнительный отступ (пробелы) добавляет вложенность внутри списка
    - обычная строка -> вложена под текущий заголовок/слайд

    Алгоритм — "стек последнего узла на каждом уровне", O(n) по строкам.
    Стек защищён от "прыжков" уровня (например, сразу с 1 на 4): если
    новый уровень больше текущей глубины стека, он не создаёт фиктивные
    промежуточные уровни, а мягко прижимается к глубине стека — дерево
    остаётся валидным при любой структуре входного текста.
    """
    root = TreeNode(name=root_name, node_type="document", level=0)
    stack: list[TreeNode] = [root]  # stack[k] = последний узел уровня k
    current_heading_level = 0  # уровень последнего заголовка/слайда

    for raw_line in text.split("\n"):
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            continue
        if skip_images and _IMAGE_LINE_RE.match(stripped):
            # строки вида ![alt](file.jpg) — подписи картинок markitdown,
            # не несут структурной информации, засоряют дерево слайдов
            continue

        slide_match = _SLIDE_MARKER_RE.match(stripped)
        if slide_match:
            level = 1
            item_text = f"Слайд {slide_match.group(1)}"
            node_type = "slide"
            current_heading_level = 1
        elif stripped.startswith("#"):
            hashes = len(stripped) - len(stripped.lstrip("#"))
            level = hashes
            item_text = stripped[hashes:].strip()
            node_type = "heading"
            current_heading_level = level
        elif stripped[:2] in ("- ", "* "):
            indent = len(line) - len(line.lstrip(" "))
            level = current_heading_level + 1 + indent // 2
            item_text = stripped[2:].strip()
            node_type = "bullet"
        else:
            indent = len(line) - len(line.lstrip(" "))
            level = current_heading_level + 1 + indent // 2
            item_text = stripped
            node_type = "line"

        # обрезаем стек до нужной глубины и, если запрошенный уровень
        # "перепрыгнул" длину стека, прижимаем его к максимально
        # доступной глубине — без этого при неровных отступах в
        # реальных документах легко получить IndexError или дыры в дереве
        while len(stack) > level:
            stack.pop()
        level = min(level, len(stack))

        node = TreeNode(name=item_text[:60], node_type=node_type, level=level, content=item_text)
        parent = stack[-1]
        parent.add_child(node)
        stack.append(node)

    return root


def parse_txt_or_md(path: str) -> TreeNode:
    """Родной .txt/.md файл — читаем и передаём в общее ядро."""
    text = pathlib.Path(path).read_text(encoding="utf-8", errors="ignore")
    return _markdown_text_to_tree(text, os.path.basename(path))


def parse_via_markitdown(path: str, skip_images: bool = True) -> TreeNode:
    """
    Универсальный парсер для 'тяжёлых' офисных форматов (docx, pptx, pdf
    и других, что понимает markitdown: xlsx, изображения с OCR и т.д.)

    Почему так лучше, чем парсить формат напрямую через python-docx/pptx:
    - markitdown — библиотека Microsoft, специально созданная, чтобы
      привести ЛЮБОЙ офисный формат к единому нормализованному markdown.
      Она уже умеет отличать таблицы, списки, слайды, изображения и
      обрабатывает массу реальных "грязных" файлов лучше, чем
      самописная эвристика по стилям абзацев.
    - Все форматы в итоге проходят через ОДНО и то же markdown-ядро
      (_markdown_text_to_tree). Значит, если мы улучшаем распознавание
      структуры (например, вложенные списки), улучшение сразу работает
      для docx, pptx И pdf одновременно — не нужно чинить три парсера.
    - ВАЖНО: markitdown восстанавливает структуру только там, где она
      реально ЕСТЬ в файле (стили заголовков, маркированные списки,
      разрывы слайдов). Если в исходном docx текст оформлен как "обычный
      абзац" без какой-либо разметки (так бывает, когда человек делает
      заголовки просто на глаз, без стиля Word) — структуры там физически
      нет, и никакой парсер её не восстановит. В этом случае дерево
      закономерно получится плоским — это отражает сам документ, а не
      ошибка парсинга.
    """
    from markitdown import MarkItDown

    converter = MarkItDown()
    result = converter.convert(path)
    return _markdown_text_to_tree(
        result.text_content, os.path.basename(path), skip_images=skip_images
    )


def parse_docx_native(path: str) -> TreeNode:
    """
    Альтернатива parse_via_markitdown: напрямую через python-docx.
    Даёт точную структуру ТОЛЬКО если в файле реально использованы
    встроенные стили 'Heading 1..9' — иначе (как в документах, где
    заголовки просто набраны обычным текстом) дерево будет плоским,
    как и через markitdown. Плюс этого варианта — отдельные узлы для
    таблиц с содержимым по ячейкам, чего markdown-конвертация не всегда
    передаёт так же явно.
    Word-документ: используем встроенные стили абзацев.
    'Heading 1' -> level 1, 'Heading 2' -> level 2 и т.д.
    Обычные абзацы становятся детьми ближайшего предыдущего заголовка.
    python-docx уже читает XML через lxml — быстро и без ручного парсинга.
    """
    import docx  # python-docx

    document = docx.Document(path)
    root = TreeNode(name=os.path.basename(path), node_type="document", level=0)
    stack: list[TreeNode] = [root]  # stack[level] = последний узел этого уровня

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style = (para.style.name or "").strip()
        if style.startswith("Heading"):
            try:
                level = int(style.split()[-1])
            except ValueError:
                level = 1
            node = TreeNode(name=text, node_type="heading", level=level)
            while len(stack) > level:
                stack.pop()
            stack[-1].add_child(node)
            stack.append(node)
        else:
            parent = stack[-1]
            node = TreeNode(
                name=text[:40] + ("..." if len(text) > 40 else ""),
                node_type="paragraph",
                level=len(stack),
                content=text,
            )
            parent.add_child(node)

    # таблицы docx тоже являются частью структуры документа
    for i, table in enumerate(document.tables, start=1):
        t_node = TreeNode(name=f"Table {i}", node_type="table", level=1)
        for r_idx, row in enumerate(table.rows):
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            t_node.add_child(
                TreeNode(name=f"Row {r_idx + 1}", node_type="row", level=2, content=row_text)
            )
        root.add_child(t_node)

    return root


def parse_pptx_native(path: str) -> TreeNode:
    """
    Альтернатива parse_via_markitdown: напрямую через python-pptx.
    Презентация: Presentation -> Slide -> Shape -> абзацы текста.
    Естественная иерархия самого формата, поэтому дерево строится
    без каких-либо эвристик — просто зеркалим объектную модель pptx.
    """
    from pptx import Presentation

    prs = Presentation(path)
    root = TreeNode(name=os.path.basename(path), node_type="document", level=0)

    for i, slide in enumerate(prs.slides, start=1):
        slide_node = TreeNode(name=f"Slide {i}", node_type="slide", level=1)
        root.add_child(slide_node)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            shape_node = TreeNode(
                name=shape.shape_type.__str__() if shape.shape_type else "Shape",
                node_type="shape",
                level=2,
            )
            slide_node.add_child(shape_node)
            for para in shape.text_frame.paragraphs:
                p_text = "".join(run.text for run in para.runs).strip()
                if p_text:
                    shape_node.add_child(
                        TreeNode(name=p_text[:40], node_type="paragraph", level=3, content=p_text)
                    )

    return root


def parse_pdf_native(path: str) -> TreeNode:
    """
    Альтернатива parse_via_markitdown: напрямую через pdfplumber.
    PDF: постраничный генератор pdfplumber — страницы обрабатываются
    последовательно, без загрузки всего документа в одну гигантскую
    структуру сразу. Внутри страницы группируем текст по строкам.
    """
    import pdfplumber

    root = TreeNode(name=os.path.basename(path), node_type="document", level=0)

    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            page_node = TreeNode(name=f"Page {i}", node_type="page", level=1)
            root.add_child(page_node)
            text = page.extract_text() or ""
            for line in text.split("\n"):
                line = line.strip()
                if line:
                    page_node.add_child(
                        TreeNode(name=line[:40], node_type="line", level=2, content=line)
                    )
    return root


def parse_json(path: str) -> TreeNode:
    """
    JSON: структура файла УЖЕ является деревом (dict/list),
    поэтому просто рекурсивно оборачиваем её в TreeNode.
    Это самый "чистый" пример универсальности модели: никакой
    эвристики, прямое отображение данных на дерево.
    """
    with open(path, encoding="utf-8") as f:
        data = json.load(f)

    root = TreeNode(name=os.path.basename(path), node_type="document", level=0)

    def build(value: Any, parent: TreeNode, level: int, key_name: str) -> None:
        if isinstance(value, dict):
            node = TreeNode(name=key_name, node_type="json_object", level=level)
            parent.add_child(node)
            for k, v in value.items():
                build(v, node, level + 1, str(k))
        elif isinstance(value, list):
            node = TreeNode(name=key_name, node_type="json_array", level=level)
            parent.add_child(node)
            for idx, item in enumerate(value):
                build(item, node, level + 1, f"[{idx}]")
        else:
            parent.add_child(
                TreeNode(name=key_name, node_type="json_value", level=level, content=str(value))
            )

    build(data, root, 1, "root")
    return root


def parse_xml_or_html(path: str) -> TreeNode:
    """
    XML/HTML: дерево тегов -> дерево TreeNode один в один.
    Используем lxml (C-реализация) — существенно быстрее, чем
    ручной парсер на регулярных выражениях, и корректно обрабатывает
    вложенность, атрибуты, самозакрывающиеся теги и т.д.
    """
    from lxml import etree

    parser = etree.HTMLParser() if path.lower().endswith((".htm", ".html")) else None
    tree = etree.parse(path, parser=parser) if parser else etree.parse(path)
    xml_root = tree.getroot()

    def build(elem, level: int) -> TreeNode:
        text = (elem.text or "").strip()
        node = TreeNode(
            name=f"<{elem.tag}>",
            node_type="xml_tag",
            level=level,
            content=text or None,
        )
        for child in elem:
            node.add_child(build(child, level + 1))
        return node

    root = TreeNode(name=os.path.basename(path), node_type="document", level=0)
    root.add_child(build(xml_root, 1))
    return root


# ---------------------------------------------------------------------------
# 3. ДИСПЕТЧЕР: выбор парсера по расширению файла
# ---------------------------------------------------------------------------

PARSERS: dict[str, Callable[[str], TreeNode]] = {
    ".txt": parse_txt_or_md,
    ".md": parse_txt_or_md,
    ".docx": parse_via_markitdown,  # по умолчанию через markitdown (см. выше)
    ".pptx": parse_via_markitdown,
    ".pdf": parse_via_markitdown,
    ".xlsx": parse_via_markitdown,
    ".json": parse_json,
    ".xml": parse_xml_or_html,
    ".html": parse_xml_or_html,
    ".htm": parse_xml_or_html,
}

# "родные" парсеры без промежуточного markdown — используйте их напрямую,
# если файл гарантированно использует стили Word/структуру pptx и вам
# нужна максимально точная структура (например, таблицы docx по ячейкам):
#   from file_tree_builder import parse_docx_native, parse_pptx_native, parse_pdf_native
#   tree = parse_docx_native("report.docx")


def build_tree(path: str) -> TreeNode:
    """
    Главная точка входа. Определяет формат по расширению файла
    и вызывает соответствующий парсер. При неизвестном расширении —
    graceful fallback: пробуем прочитать как обычный текст.
    """
    if not os.path.isfile(path):
        raise FileNotFoundError(f"Файл не найден: {path}")

    ext = os.path.splitext(path)[1].lower()
    parser = PARSERS.get(ext)

    if parser is None:
        # универсальный запасной путь — если формат неизвестен,
        # но файл текстовый, всё равно строим дерево по строкам
        try:
            return parse_txt_or_md(path)
        except UnicodeDecodeError as exc:
            raise ValueError(
                f"Неизвестный или бинарный формат '{ext}', добавьте свой парсер в словарь PARSERS."
            ) from exc

    try:
        return parser(path)
    except ImportError as exc:
        raise ImportError(
            f"Для чтения '{ext}' нужна библиотека: {exc}. "
            f"Установите её командой pip install <имя-пакета>."
        ) from exc


# ---------------------------------------------------------------------------
# 4. ПРИМЕР ИСПОЛЬЗОВАНИЯ
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    file_path = "Курсовая 8 вариант Дремин.docx"
    tree = build_tree(file_path)

    print(f"\nДерево файла: {file_path}")
    print(f"Всего узлов: {tree.count_nodes()}\n")
    tree.print_tree()

    # экспорт в JSON — работает одинаково для ЛЮБОГО формата исходника,
    # потому что структура TreeNode всегда одна и та же.
    # Пишем рядом со скриптом (а не рядом с исходником), т.к. папка
    # с исходным файлом может быть доступна только для чтения.
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    out_path = os.path.join(os.getcwd(), base_name + "_tree.json")
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(tree.to_dict(), f, ensure_ascii=False, indent=2)
    print(f"\nJSON-дерево сохранено в: {out_path}")
